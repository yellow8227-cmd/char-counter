# -*- coding: utf-8 -*-
"""
赫血 KAKKETSU — 「그것」(封神) 캐릭터 컨셉 레퍼런스
한국 투자사 제출용 비주얼 덱 빌더.

    python3 build_deck_kr.py

assets/ 의 이미지를 읽어 build/KAKKETSU_그것_컨셉레퍼런스_KR.pptx 를 만든다.
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
A = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "build", "KAKKETSU_그것_컨셉레퍼런스_KR.pptx")

SW, SH = 13.3333, 7.5          # slide size (in)
ML, MR = 0.78, 0.78            # side margins
TOP = 0.52

# ── palette ────────────────────────────────────────────────────────────────
INK      = RGBColor(0x0A, 0x08, 0x09)   # near black
BONE     = RGBColor(0xEC, 0xE6, 0xDA)   # primary text on dark
BONE_MID = RGBColor(0xB7, 0xAD, 0x9F)
BONE_DIM = RGBColor(0x8A, 0x81, 0x76)
RED      = RGBColor(0xB4, 0x23, 0x1A)   # accent — 부적의 붉은색
RED_DIM  = RGBColor(0x8A, 0x44, 0x3B)
DARKINK  = RGBColor(0x17, 0x14, 0x14)   # text on light slides
DARKMID  = RGBColor(0x4B, 0x44, 0x40)

KR = "Malgun Gothic"     # 한글 본문/제목
LAT = "Cambria"          # 라틴 소제목·숫자

PALETTE = [
    ("0F0E0E", "칠흑"),      ("353434", "재"),
    ("4B4844", "마른 회"),   ("6F6357", "젖은 흙"),
    ("C4B6A6", "삼베·뼈"),   ("A9635C", "빛바랜 홍"),
    ("5F686A", "색동의 청"), ("8A443B", "핏빛"),
]

pageno = {"n": 0}


# ── low-level helpers ──────────────────────────────────────────────────────
def set_alpha(shape, pct):
    """솔리드 채움에 투명도(%) 부여."""
    spPr = shape.fill._xPr
    solid = spPr.find(qn("a:solidFill"))
    srgb = solid.find(qn("a:srgbClr"))
    srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (nsdecls("a"), int(pct * 1000))))


def set_line_alpha(shape, pct):
    ln = shape.line._get_or_add_ln()
    solid = ln.find(qn("a:solidFill"))
    srgb = solid.find(qn("a:srgbClr"))
    srgb.append(parse_xml('<a:alpha %s val="%d"/>' % (nsdecls("a"), int(pct * 1000))))


def cjk(run, face):
    """한글/한자가 지정 폰트로 렌더되도록 ea·cs 도 함께 설정."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = parse_xml("<%s %s/>" % (tag, nsdecls("a")))
            rPr.append(el)
        el.set("typeface", face)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, *, font=KR, bold=False, spc=0, line=1.35,
         before=0, after=0, align=PP_ALIGN.LEFT, first=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
    pPr.set("lnSpc", "") if False else None
    p.line_spacing = line
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    if spc:
        r._r.get_or_add_rPr().set("spc", str(int(spc * 100)))
    cjk(r, font if font != LAT else KR)
    return p


def pic_cover(slide, name, x, y, w, h, *, focus=0.5, behind=True):
    """지정 박스를 꽉 채우도록 이미지를 crop-to-fill 배치."""
    path = os.path.join(A, name + ".jpg")
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    shp = slide.shapes.add_picture(path, In(x), In(y), In(w), In(h))
    if img_ar > box_ar:                      # 이미지가 더 넓다 → 좌우 crop
        keep = box_ar / img_ar
        side = (1 - keep)
        shp.crop_left = side * focus
        shp.crop_right = side * (1 - focus)
    else:                                    # 이미지가 더 높다 → 상하 crop
        keep = img_ar / box_ar
        side = (1 - keep)
        shp.crop_top = side * focus
        shp.crop_bottom = side * (1 - focus)
    if behind:
        sp = shp._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)
    return shp


def pic_fit(slide, name, x, y, w, h):
    """박스 안에 이미지 전체가 들어가도록(letterbox) 중앙 배치. 실제 rect 반환."""
    path = os.path.join(A, name + ".jpg")
    iw, ih = Image.open(path).size
    ar = iw / ih
    if ar > w / h:
        nw, nh = w, w / ar
    else:
        nh, nw = h, h * ar
    nx, ny = x + (w - nw) / 2, y + (h - nh) / 2
    slide.shapes.add_picture(path, In(nx), In(ny), In(nw), In(nh))
    return nx, ny, nw, nh


def scrim(slide, name, x=0, y=0, w=SW, h=SH):
    slide.shapes.add_picture(os.path.join(A, name + ".png"), In(x), In(y), In(w), In(h))


def rect(slide, x, y, w, h, rgb, alpha=100, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, In(x), In(y), In(w), In(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if alpha < 100:
        set_alpha(s, alpha)
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def hairline(slide, x, y, w, rgb=BONE_DIM, alpha=45, weight=0.6):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(x), In(y), In(w), Pt(weight))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    set_alpha(s, alpha)
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def blank(prs):
    """검정 배경 슬라이드. 배경은 p:bg 로 칠해 도형 z-order 를 건드리지 않는다."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    cSld = s._element.find(qn("p:cSld"))
    bg = parse_xml(
        '<p:bg %s><p:bgPr><a:solidFill><a:srgbClr val="0A0809"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
        % nsdecls("p", "a"))
    cSld.insert(0, bg)
    return s


# ── recurring furniture ────────────────────────────────────────────────────
def head(slide, index, eyebrow, title, *, y=TOP, color=BONE, dim=BONE_DIM,
         accent=RED, sub=None, width=None):
    tf = textbox(slide, ML, y, width or (SW - ML - MR), 0.32)
    p = para(tf, "%s" % index, 11.5, accent, font=LAT, bold=True, spc=1.4,
             line=1.0, first=True)
    r = p.add_run()
    r.text = "   " + eyebrow
    r.font.size, r.font.bold, r.font.name = Pt(9.5), True, LAT
    r.font.color.rgb = dim
    r._r.get_or_add_rPr().set("spc", "180")
    cjk(r, KR)

    tf2 = textbox(slide, ML, y + 0.36, width or (SW - ML - MR), 0.8)
    para(tf2, title, 33, color, bold=True, line=1.12, spc=-0.6, first=True)
    if sub:
        tf3 = textbox(slide, ML, y + 1.12, width or (SW - ML - MR), 0.4)
        para(tf3, sub, 13, dim, line=1.35, first=True)


def foot(slide, *, light=False, x=ML):
    pageno["n"] += 1
    c = DARKMID if light else BONE_DIM
    tf = textbox(slide, x, SH - 0.52, 7.0, 0.24)
    p = para(tf, "赫血 KAKKETSU   ·   「그것」封神 캐릭터 컨셉 레퍼런스", 7.5, c,
             spc=1.2, line=1.0, first=True)
    if not light:
        p.runs[0].font.color.rgb = RGBColor(0x6E, 0x66, 0x5D)
    tf2 = textbox(slide, SW - MR - 1.6, SH - 0.52, 1.6, 0.24)
    para(tf2, "%02d" % pageno["n"], 8.5, c, font=LAT, spc=1.0, line=1.0,
         align=PP_ALIGN.RIGHT, first=True)


def note_card(slide, x, y, w, label, body, *, size=10.5, lab_size=8.5,
              pad=0.24, alpha=62, accent=RED, minh=0.9, body_color=None):
    """주석 카드 — 반투명 검정 판 + 라벨 + 본문."""
    cpl = max(8.0, (w - pad * 2) * 72.0 / size * 1.06)   # 한 줄에 들어가는 한글 자수
    lines = 0
    for chunk in body.split("\n"):
        lines += max(1, -(-len(chunk) // int(cpl)))
    h = max(minh, pad * 2 + 0.20 + lines * size * 1.46 / 72.0)
    plate = rect(slide, x, y, w, h, RGBColor(0x08, 0x06, 0x07), alpha)
    tf = textbox(slide, x + pad, y + pad, w - pad * 2, h - pad * 2)
    para(tf, label, lab_size, accent, bold=True, spc=1.6, line=1.0, first=True)
    para(tf, body, size, body_color or BONE_MID, line=1.5, before=5)
    return h


def callout(slide, tx, ty, lx, ly, lw, label, body="", *, side="left",
            lab_size=9.5, body_size=8.5):
    """이미지 위 지시선 주석: 표적점(tx,ty) — 라벨 박스(lx,ly)."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(tx - 0.045), In(ty - 0.045),
                                 In(0.09), In(0.09))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED
    dot.line.color.rgb = BONE; dot.line.width = Pt(0.75)
    set_line_alpha(dot, 70)
    dot.shadow.inherit = False

    x0, x1 = (lx + lw, tx) if side == "left" else (tx, lx)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(min(x0, x1)),
                                In(ty - 0.004), In(abs(x1 - x0)), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = BONE
    set_alpha(ln, 42)
    ln.line.fill.background(); ln.shadow.inherit = False

    align = PP_ALIGN.RIGHT if side == "left" else PP_ALIGN.LEFT
    tf = textbox(slide, lx, ly, lw, 0.7)
    para(tf, label, lab_size, BONE, bold=True, line=1.15, align=align, first=True)
    if body:
        para(tf, body, body_size, BONE_MID, line=1.4, before=2, align=align)


# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = In(SW), In(SH)


# ── 01 · COVER ─────────────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "corridor_hospital", 0, 0, SW, SH, focus=0.5)
scrim(s, "vignette")
scrim(s, "scrim_left", 0, 0, SW * 0.72, SH)
scrim(s, "scrim_bottom", 0, SH * 0.45, SW, SH * 0.55)

tf = textbox(s, ML, 0.62, 6.0, 0.3)
para(tf, "WALLAH PROJECT   ×   SUPERNOVA STUDIOS", 9, BONE_MID, font=LAT,
     bold=True, spc=2.2, line=1.0, first=True)

s.shapes.add_picture(os.path.join(A, "logo_kakketsu.png"), In(ML - 0.42), In(2.20),
                     In(4.8), In(1.19))

tf = textbox(s, ML, 1.72, 8.0, 0.4)
para(tf, "CHARACTER CONCEPT REFERENCE", 10, RED, font=LAT, bold=True,
     spc=3.4, line=1.0, first=True)

tf = textbox(s, ML, 3.62, 8.6, 1.9)
para(tf, "각혈  ·  KAKKETSU", 15, BONE_MID, spc=2.0, line=1.2, first=True)
para(tf, "「그것」— 封神 (봉신)", 42, BONE, bold=True, line=1.15, spc=-0.8, before=10)
para(tf, "저주라 믿었던 존재의 설계도", 14.5, BONE_MID, line=1.4, before=12)

hairline(s, ML, 5.86, 2.3, BONE, 34)

tf = textbox(s, ML, 6.06, 8.0, 0.8)
para(tf, "각본 · 감독   한동하   |   장편 · 공포 · 오컬트   |   한국 – 일본", 11,
     BONE_MID, line=1.5, first=True)
para(tf, "2026 투자 검토용 캐릭터 비주얼 자료", 10, RED_DIM, line=1.4, before=4)

tf = textbox(s, SW - MR - 4.2, SH - 0.78, 4.2, 0.3)
para(tf, "CONFIDENTIAL  ·  FOR INVESTMENT REVIEW ONLY", 8, RGBColor(0x6E, 0x66, 0x5D),
     font=LAT, spc=1.6, line=1.0, align=PP_ALIGN.RIGHT, first=True)


# ── 02 · LOGLINE (light) ───────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "fog_white", 0, 0, SW, SH, focus=0.24)
scrim(s, "scrim_light_left", 0, 0, SW, SH)

tf = textbox(s, ML, 1.15, 3.0, 0.3)
para(tf, "LOGLINE", 10.5, RED, font=LAT, bold=True, spc=3.2, line=1.0, first=True)

tf = textbox(s, ML, 1.72, 6.55, 3.4)
para(tf, "대대로 이어진 가문의 저주를 끊기 위해\n싸워온 한 남자는, 그 저주라 믿었던 존재가\n사실은 가족을 지켜온 수호였음을 깨닫고\n핏줄과 죄의 진짜 의미와 마주한다.",
     22, DARKINK, bold=True, line=1.62, spc=-0.4, first=True)

hairline(s, ML, 5.42, 1.7, DARKMID, 55)

tf = textbox(s, ML, 5.62, 6.2, 1.0)
para(tf, "4대에 걸친 핏줄, 하나의 형상. 「그것」은 이 이야기의 괴물이자,\n이 이야기의 마지막 진실이다.",
     11.5, DARKMID, line=1.6, first=True)
foot(s, light=True)


# ── 03 · 기획의도 ──────────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "corridor_wood", 0, 0, SW, SH, focus=0.45)
scrim(s, "vignette")
scrim(s, "scrim_left", 0, 0, SW * 0.80, SH)
head(s, "01", "DIRECTOR'S INTENTION", "기획의도")

tf = textbox(s, ML, 2.05, 6.15, 2.4)
para(tf, "이 작품은 공포영화의 문법을 빌려,\n‘지켜보는 존재’에 대한 이야기를 한다.",
     19, BONE, bold=True, line=1.5, spc=-0.4, first=True)
para(tf, "쫓아오고, 따라붙고, 침묵하는 존재. 우리는 그것을 언제나 저주라 불러왔다.\n"
         "그러나 이 영화는 묻는다. 그것이 정말 파괴를 위해 존재했는가, 아니면\n"
         "우리가 외면해 온 죄를 대신 짊어진 결과였는가.",
     11.5, BONE_MID, line=1.7, before=16)

note_card(s, ML, 5.05, 6.15,
          "관객이 가져가는 감정",
          "이 영화가 끝났을 때, 관객은 더 이상 ‘귀신’을 무서워하지 않게 될 것이다.\n"
          "대신, 인간의 선택을 두려워하게 될 것이다.",
          size=11.5, alpha=70)
foot(s)


# ── 04 · THE TURN ──────────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "kneel_closeup", SW * 0.46, 0, SW * 0.54, SH, focus=0.42)
rect(s, 0, 0, SW * 0.50, SH, INK)
scrim(s, "scrim_left", SW * 0.44, 0, SW * 0.30, SH)
head(s, "02", "THE TURN · 반전 구조", "저주가 아니라,\n수호였다", width=5.6)

tf = textbox(s, ML, 2.26, 5.57, 0.7)
para(tf, "「그것」은 저주가 아니었다. 이 집안이 저질러 온 죄와 선택의 대가를\n"
         "스스로 떠안고, 끝까지 막아 서 있던 하나의 수호였다.",
     11.5, BONE, bold=True, line=1.55, first=True)
hairline(s, ML, 3.00, 5.57, BONE, 22)

steps = [
    ("01", "따라붙는다", "1582년 조선의 처형장에서 태어나 4대에 걸쳐 이동한다.\n이름도 얼굴도 기억하지 않는다. 기준은 오직 하나 — 핏줄."),
    ("02", "죽어 나간다", "할아버지 동필, 아버지 종문. 두 세대가 같은 자리에서\n같은 방식으로 사라진다. 저주는 다음 대로 넘어간다."),
    ("03", "사라진다", "아들이 죽고 「그것」은 걸음을 되찾아 문을 나선다.\n그리고 그때, 비로소 진짜 지옥이 시작된다."),
]
y = 3.34
for i, (n, t, b) in enumerate(steps):
    tfn = textbox(s, ML, y - 0.04, 0.6, 0.4)
    para(tfn, n, 15, RED, font=LAT, bold=True, line=1.0, first=True)
    tf = textbox(s, ML + 0.62, y - 0.06, 4.95, 1.1)
    para(tf, t, 14.5, BONE, bold=True, line=1.2, first=True)
    para(tf, b, 9.8, BONE_MID, line=1.55, before=5)
    if i < 2:
        hairline(s, ML, y + 0.98, 5.57, BONE, 18)
    y += 1.16
foot(s)


# ── 05 · 「그것」 정의 ──────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "hall_tilt", 0, 0, SW, SH, focus=0.5)
scrim(s, "vignette")
scrim(s, "scrim_bottom", 0, SH * 0.28, SW, SH * 0.72)
scrim(s, "scrim_top", 0, 0, SW, SH * 0.42)
scrim(s, "scrim_bottom", 0, SH * 0.62, SW, SH * 0.38)
head(s, "03", "WHAT IS 「IT」", "「그것」")

tf = textbox(s, ML, 4.28, 7.4, 1.2)
para(tf, "「그것」은 특정 개인이나 유령이 아니다.\n"
         "인간의 죄와 피, 그리고 말해지지 못한 원념이 쌓이고 쌓여,\n"
         "마침내 형태를 얻은 결과로서 태어난 존재이다.",
     15, BONE, bold=True, line=1.6, spc=-0.3, first=True)

chips = [
    ("탄생", "1587년 조선 평양 — 무당의 굿판"),
    ("형상", "노모의 시신 + 「封神」 부적"),
    ("규칙", "이름도 얼굴도 아닌, 오직 핏줄"),
    ("정체", "저주의 형상을 한 봉인(封)"),
]
x = ML
for k, v in chips:
    tf = textbox(s, x, 6.16, 3.0, 0.6)
    para(tf, k, 9, RED, bold=True, spc=1.8, line=1.0, first=True)
    para(tf, v, 10.5, BONE_MID, line=1.35, before=5)
    x += 3.0
foot(s)


# ── 06 · 캐릭터 스펙 ───────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "front_stand", SW * 0.415, 0, SW * 0.585, SH, focus=0.5)
rect(s, 0, 0, SW * 0.46, SH, INK)
scrim(s, "scrim_left", SW * 0.40, 0, SW * 0.26, SH)
head(s, "04", "CHARACTER SPEC", "封神  봉신", width=5.0)

spec = [("키", "약 150cm  (연출 스케일 180 / 190 / 260cm)"),
        ("성별", "여성"),
        ("나이", "노인"),
        ("정체", "망나니의 노모 — 굿판에서 봉인된 시신")]
y = 2.12
for k, v in spec:
    tf = textbox(s, ML, y, 0.62, 0.26)
    para(tf, k, 9, RED, bold=True, spc=1.2, line=1.0, first=True)
    tf = textbox(s, ML + 0.72, y - 0.02, 4.3, 0.3)
    para(tf, v, 11, BONE, line=1.3, first=True)
    hairline(s, ML, y + 0.30, 4.95, BONE, 16)
    y += 0.50

tf = textbox(s, ML, 4.30, 4.95, 2.3)
para(tf, "특징", 9, RED, bold=True, spc=1.8, line=1.0, first=True)
para(tf, "이마의 붉은 부적으로 인해 얼굴을 볼 수 없다.\n"
         "매우 마른 체형, 메마른 피부. 빛바랜 색동저고리의 누더기 차림.\n"
         "한쪽 다리를 절룩이며 무겁게 이동한다.\n\n"
         "등 뒤에 거대한 원혼 덩어리가 단단히 붙어 있고, 긴 머리카락이\n"
         "그 전체를 감싸 원혼 덩어리를 억누르고 있다.",
     10, BONE_MID, line=1.62, before=7)

callout(s, 9.35, 3.02, 6.28, 2.86, 2.50, "이마의 부적 「封神」", "얼굴이 있어야 할 자리")
callout(s, 9.42, 4.44, 6.28, 4.28, 2.50, "빛바랜 색동저고리", "색이 거의 빠진 누더기")
callout(s, 10.72, 1.82, 11.02, 1.66, 1.53, "원혼 덩어리", "등 뒤에 업힌 수백의 원혼",
        side="right")
callout(s, 9.30, 6.52, 6.28, 6.36, 2.50, "뒤틀린 다리 · 짚신", "질질 끌리는 한쪽 발")
foot(s)


# ── 07 · 디자인 앵커 ① 부적 ────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "front_arms", SW * 0.30, 0, SW * 0.70, SH, focus=0.5)
rect(s, 0, 0, SW * 0.34, SH, INK)
scrim(s, "scrim_left", SW * 0.28, 0, SW * 0.30, SH)
head(s, "05", "DESIGN ANCHOR 01", "부적 「封神」\n— 얼굴을 대신하는 봉인", width=5.3)

tf = textbox(s, ML, 2.62, 4.55, 1.4)
para(tf, "한국의 부적은 일본의 오마모리처럼 ‘좋은 기운을 불러들이는 것’이\n"
         "아니라, 위험한 존재나 힘을 눌러 봉해두기 위한 것에 가깝다.\n"
         "일본적 감각으로는 후다(札)·인(印)·결계(結界)에 가깝다.",
     10, BONE_MID, line=1.62, first=True)

note_card(s, ML, 4.30, 4.55, "연출 노트",
          "얼굴이 없다 = 감정을 읽을 수 없다. 관객은 끝까지 「그것」의 의도를\n"
          "판독하지 못하고, 그 판독 불가능성이 이 영화의 서스펜스를 지탱한다.\n\n"
          "그리고 마지막, 슈가 머리카락을 걷어낸 자리에는 — 아무것도 없다.",
          size=9.8, alpha=66)

s.shapes.add_picture(os.path.join(A, "detail_talisman_face.jpg"),
                     In(ML), In(5.86), In(0.55), In(0.850))
tf = textbox(s, ML + 0.73, 5.96, 3.7, 0.6)
para(tf, "「封神」 — 봉할 봉, 귀신 신", 10, BONE, bold=True, line=1.2, first=True)
para(tf, "귀신을 부르는 글자가 아니라, 가두는 글자다.", 9, BONE_DIM, line=1.35, before=3)
foot(s)


# ── 08 · 디자인 앵커 ② 색동저고리 ──────────────────────────────────────────
s = blank(prs)
pic_cover(s, "saekdong_back", 0, 0, SW * 0.62, SH, focus=0.45)
rect(s, SW * 0.575, 0, SW * 0.425, SH, INK)
scrim(s, "scrim_right", SW * 0.34, 0, SW * 0.28, SH)
head(s, "06", "DESIGN ANCHOR 02", "색동저고리\n— 아이를 지키는 옷",
     width=4.6)
# 우측 컬럼으로 헤더 이동
for shp in list(s.shapes)[-3:]:
    shp.left = In(SW * 0.60)

tf = textbox(s, SW * 0.60, 2.62, 4.55, 1.4)
para(tf, "색동저고리는 한국에서 아이를 지키기 위해 입히는 옷이다.\n"
         "나쁜 기운을 멀리하고 생명이 무사히 자라기를 바라는,\n"
         "주술적 보호의 의미가 담겨 있다.",
     10, BONE_MID, line=1.62, first=True)

note_card(s, SW * 0.60, 4.24, 4.55, "연출 노트",
          "가장 무서운 형상에 가장 다정한 옷을 입힌다.\n"
          "이 모순이 곧 반전의 복선이다. 색이 다 빠진 누더기는\n"
          "「그것」이 얼마나 오래 지켜왔는지를 말하는 유일한 물증이다.",
          size=9.8, alpha=66)

s.shapes.add_picture(os.path.join(A, "detail_saekdong.jpg"),
                     In(SW * 0.60), In(5.86), In(0.548), In(0.850))
tf = textbox(s, SW * 0.60 + 0.72, 5.96, 3.8, 0.6)
para(tf, "청 · 적 · 황 · 백 · 흑 — 오방색", 10, BONE, bold=True, line=1.2, first=True)
para(tf, "다섯 방위의 색으로 액을 막는다는 믿음에서 온 배색.", 9, BONE_DIM,
     line=1.35, before=3)
foot(s, x=SW * 0.60)


# ── 09 · 디자인 앵커 ③ 짚신 ────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "foot_closeup", SW * 0.50, 0, SW * 0.50, SH, focus=0.5)
rect(s, 0, 0, SW * 0.54, SH, INK)
scrim(s, "scrim_left", SW * 0.46, 0, SW * 0.26, SH)
head(s, "07", "DESIGN ANCHOR 03", "짚신과 뒤틀린 다리\n— 업을 지고 걷는 몸", width=5.6)

tf = textbox(s, ML, 2.62, 5.5, 1.2)
para(tf, "한국의 짚신은 일본의 와라지와 형상은 비슷하지만, 단순한 여행용\n"
         "신발이 아니다. 고통과 업을 짊어진 채 계속 걸어가는 인간의\n"
         "상징으로 쓰인다.",
     10, BONE_MID, line=1.62, first=True)

note_card(s, ML, 4.02, 5.5, "연출 노트 — 절룩임은 상처가 아니라 하중이다",
          "등 뒤 수백의 원혼을 대신 짊어졌기에 몸이 무거워졌고, 그래서 절 수밖에\n"
          "없었다. 지켜야 할 핏줄이 끊어져 책임이 끝나는 순간, 걸음은 정상으로\n"
          "돌아온다 — 관객이 반전을 몸으로 알아채는 지점.",
          size=9.8, alpha=66)

s.shapes.add_picture(os.path.join(A, "detail_shoe.jpg"), In(ML), In(5.86), In(2.6), In(0.885))
tf = textbox(s, ML + 2.80, 5.92, 2.7, 0.8)
para(tf, "슥……  쿵……", 15, RED, bold=True, line=1.15, first=True)
para(tf, "짚신이 바닥을 긁는 소리.\n이 캐릭터의 사운드 테마이자, 등장 신호.",
     9, BONE_DIM, line=1.4, before=4)
foot(s)


# ── 10 · 원혼 덩어리 ───────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "hall_sweep", 0, 0, SW, SH, focus=0.5)
scrim(s, "vignette")
scrim(s, "scrim_bottom", 0, SH * 0.30, SW, SH * 0.70)
scrim(s, "scrim_top", 0, 0, SW, SH * 0.40)
head(s, "08", "THE BURDEN", "등에 업힌 수백의 원혼")

items = [
    ("유동성", "액체처럼 흐른다. 고정된 실루엣이 없다."),
    ("구성", "수많은 사람의 얼굴 · 팔다리 · 손등이 얽혀 유기적으로 변형."),
    ("구속", "긴 머리카락이 덩어리 전체를 감아 억누르고 있다."),
    ("이동", "덩어리도 바닥으로 흘러내리며 함께 움직인다."),
]
x = ML
for k, v in items:
    tf = textbox(s, x, 5.28, 2.85, 1.0)
    para(tf, k, 9, RED, bold=True, spc=1.8, line=1.0, first=True)
    para(tf, v, 10, BONE_MID, line=1.5, before=6)
    x += 3.0

tf = textbox(s, ML, 6.52, 11.0, 0.4)
para(tf, "연출 포인트 — 계단을 도는 순간 「그것」이 온몸의 원혼을 떨쳐낸다. "
         "짐을 버린 몸은 폭발적으로 빨라진다. 추격의 가속 구간.",
     10.5, BONE, line=1.4, first=True)
foot(s)


# ── 11 · 룩 디벨롭먼트 A · B ───────────────────────────────────────────────
s = blank(prs)
head(s, "09", "LOOK DEVELOPMENT", "룩 디벨롭먼트  A · B")
x0, y0, w0 = ML, 1.62, (SW - ML - MR - 0.34) / 2
pic_fit(s, "sheet_wet", x0, y0, w0, 3.86)
pic_fit(s, "sheet_wet2", x0 + w0 + 0.34, y0, w0, 3.86)
for i, (t, b) in enumerate([
    ("A  ·  WET / BLACK",
     "젖은 흑색 점액 톤. 원혼 덩어리와 머리카락의 경계가 사라져\n"
     "‘하나의 검은 질량’으로 읽힌다. 어두운 세트에서 가장 강한 실루엣."),
    ("B  ·  WET / DETAIL",
     "동일 톤에서 개체 식별성을 끌어올린 버전. 원혼 하나하나의 얼굴이\n"
     "읽히므로 클로즈업과 반전 시퀀스에 유리하다.")]):
    xx = x0 + i * (w0 + 0.34)
    tf = textbox(s, xx, 5.62, w0, 0.9)
    para(tf, t, 12, BONE, bold=True, spc=0.8, line=1.2, first=True)
    para(tf, b, 9.5, BONE_MID, line=1.55, before=6)
foot(s)


# ── 12 · 룩 디벨롭먼트 C · D ───────────────────────────────────────────────
s = blank(prs)
head(s, "10", "LOOK DEVELOPMENT", "룩 디벨롭먼트  C · D")
pic_fit(s, "sheet_color", x0, y0, w0, 3.86)
pic_fit(s, "sheet_mist", x0 + w0 + 0.34, y0, w0, 3.86)
for i, (t, b) in enumerate([
    ("C  ·  SAEKDONG",
     "색동의 채도를 살린 버전. ‘아이를 지키는 옷’이라는 반전의 단서가\n"
     "가장 먼저 읽힌다. 본편 기준안으로 제안."),
    ("D  ·  MIST",
     "안개 · 저채도. 형태를 지우고 실루엣만 남긴다. 티저 · 포스터 ·\n"
     "선공개 소재로 유리한 방향.")]):
    xx = x0 + i * (w0 + 0.34)
    tf = textbox(s, xx, 5.62, w0, 0.9)
    para(tf, t, 12, BONE, bold=True, spc=0.8, line=1.2, first=True)
    para(tf, b, 9.5, BONE_MID, line=1.55, before=6)

tf = textbox(s, ML, 6.52, 11.0, 0.35)
para(tf, "제안 — 본편 C(색동) / 티저 · 키비주얼 D(안개) / 클로즈업 B(디테일) 로 분리 운용",
     10, RED_DIM, bold=True, line=1.3, first=True)
foot(s)


# ── 13 · 3면도 & 디테일 ────────────────────────────────────────────────────
s = blank(prs)
head(s, "11", "TURNAROUND & DETAIL", "3면도 · 디테일 시트")
pic_fit(s, "sheet_size", ML, 1.60, 8.35, 5.05)

tf = textbox(s, ML + 8.72, 1.66, 2.85, 0.4)
para(tf, "파트별 확정 항목", 9, RED, bold=True, spc=1.8, line=1.0, first=True)
parts = [("정면 · 측면 · 후면", "실루엣 기준 3면 확정"),
         ("이마의 부적", "위치 · 각도 · 붉기"),
         ("색동저고리 (누더기)", "탈색 단계와 찢김 패턴"),
         ("메마른 피부와 손", "혈관 · 관절 돌출"),
         ("젖은 발목", "흘러내리는 검은 액체"),
         ("원혼 덩어리", "밀도와 얼굴 노출 비율"),
         ("액체 질감의 손등", "머리카락과의 접합부")]
yy = 2.06
for k, v in parts:
    tf = textbox(s, ML + 8.72, yy, 2.85, 0.5)
    para(tf, k, 10, BONE, bold=True, line=1.2, first=True)
    para(tf, v, 8.8, BONE_DIM, line=1.3, before=2)
    yy += 0.58

tf = textbox(s, ML + 8.72, 6.28, 2.85, 0.5)
para(tf, "150cm · 180cm · 190cm 세 크기에서\n동일 디자인의 일관성을 검증한 시트.",
     9, BONE_DIM, line=1.45, first=True)
foot(s)


# ── 14 · 스케일 ────────────────────────────────────────────────────────────
s = blank(prs)
head(s, "12", "CHARACTER SCALE", "크기 설계 — 같은 존재, 다른 압력")
pic_fit(s, "scale_chart", ML, 1.56, 8.55, 4.35)
pic_fit(s, "sheet_montage", ML + 8.75, 1.56, 2.65, 4.35)

rows = [("150cm", "일상 · 미행 — 사람 크기. 무해해 보여야 한다."),
        ("180cm", "대면 — 인물과 눈높이가 맞는 순간의 압박."),
        ("190cm", "추격 — 공간을 좁히는 크기."),
        ("260cm", "굿판 · 강림 — 공간 자체를 압도한다.")]
x = ML
for k, v in rows:
    tf = textbox(s, x, 6.14, 2.72, 0.8)
    para(tf, k, 12, RED, font=LAT, bold=True, line=1.0, first=True)
    para(tf, v, 9.3, BONE_MID, line=1.45, before=5)
    x += 2.87
foot(s)


# ── 15 · 의상 컨셉 ─────────────────────────────────────────────────────────
s = blank(prs)
head(s, "13", "COSTUME CONCEPT", "의상 컨셉 — 저고리 & 치마")
cw = (SW - ML - MR - 0.56) / 3
for i, n in enumerate(["cost_skirt_a", "cost_skirt_b", "cost_skirt_c"]):
    pic_fit(s, n, ML + i * (cw + 0.28), 1.62, cw, 3.9)
labs = [("변형 A", "저채도 · 흙빛. 세트에 묻히는 톤."),
        ("변형 B", "색동 채도 유지. 반전 단서 가시성 ↑"),
        ("변형 C", "삭아 풀린 올. 세월과 하중의 질감.")]
for i, (t, b) in enumerate(labs):
    tf = textbox(s, ML + i * (cw + 0.28), 5.68, cw, 0.8)
    para(tf, t, 11, BONE, bold=True, line=1.2, first=True)
    para(tf, b, 9.3, BONE_MID, line=1.45, before=5)

tf = textbox(s, ML, 6.62, 11.4, 0.4)
para(tf, "구성 — 색동 소매 저고리 + 붉은 끈 · 술 노리개 + 갈래갈래 찢긴 치마. "
         "동일 패턴을 저고리/치마/바지 3벌로 전개.",
     10, BONE_MID, line=1.4, first=True)
foot(s)


# ── 16 · 의상 변형 (바지) ──────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "cost_pants", SW * 0.40, 0, SW * 0.60, SH, focus=0.5)
rect(s, 0, 0, SW * 0.44, SH, INK)
scrim(s, "scrim_left", SW * 0.36, 0, SW * 0.28, SH)
head(s, "14", "COSTUME VARIANT", "의상 변형\n— 저고리 & 바지", width=5.0)

tf = textbox(s, ML, 2.70, 4.6, 1.2)
para(tf, "액션 · 추격 · 와이어 촬영용 변형. 치마 실루엣을 유지하지 못하는\n"
         "장면(계단 추격, 천장 매달림, 굿판 폭주)에 사용한다.",
     10, BONE_MID, line=1.6, first=True)

note_card(s, ML, 3.86, 4.6, "제작 노트",
          "붉은 바지는 색동의 ‘적(赤)’을 그대로 확대한 것이다.\n"
          "몸을 뒤집거나 거꾸로 매달았을 때에도 캐릭터의 색 기호가 유지된다.",
          size=9.8, alpha=66)

s.shapes.add_picture(os.path.join(A, "detail_talisman_tag.jpg"),
                     In(ML), In(5.56), In(0.98), In(0.861))
tf = textbox(s, ML + 1.18, 5.66, 3.4, 0.7)
para(tf, "소품 — 「封神」 부적", 10.5, BONE, bold=True, line=1.2, first=True)
para(tf, "무당이 자신의 피로 그린 부적. 극 중 유일하게\n색을 잃지 않는 소품이다.",
     9.3, BONE_DIM, line=1.45, before=4)
foot(s)


# ── 17 · 연출 ① 강림 ───────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "hall_open", 0, 0, SW, SH, focus=0.5)
scrim(s, "vignette")
scrim(s, "scrim_top", 0, 0, SW, SH * 0.40)
scrim(s, "scrim_bottom", 0, SH * 0.55, SW, SH * 0.45)
head(s, "15", "SET PIECE 01", "굿판 — 강림")

tf = textbox(s, ML, 5.62, 6.6, 1.0)
para(tf, "1587년 조선 평양. 무당은 자신의 피로 부적을 만들고,\n"
         "망나니는 노모의 무덤을 파헤친다. 부적이 이마에 붙는 순간\n"
         "시신은 인간의 형상을 잃는다.",
     11, BONE, line=1.6, first=True)

tf = textbox(s, SW - MR - 4.2, 5.62, 4.2, 1.0)
para(tf, "“그것은 괴물이 아니다.\n핏줄을 지키기 위해 모든 원혼을\n스스로 떠안은 하나의 봉인이다.”",
     11.5, BONE_MID, line=1.6, align=PP_ALIGN.RIGHT, first=True)
foot(s)


# ── 18 · 연출 ② 대청 시퀀스 ────────────────────────────────────────────────
s = blank(prs)
head(s, "16", "SET PIECE 02", "대청 — 원혼의 전개")
gw = (SW - ML - MR - 0.30) / 2
pic_fit(s, "hall_descend_a", ML, 1.66, gw, 3.32)
pic_fit(s, "hall_descend_b", ML + gw + 0.30, 1.66, gw, 3.32)
for i, (t, b) in enumerate([
    ("BEFORE  ·  응축",
     "천장에서 한 줄기로 쏟아져 내린다. 색동 소매만 색을 가진다."),
    ("AFTER  ·  전개",
     "공간 전체로 부챗살처럼 벌어진다. 얼굴들이 한꺼번에 드러나는 컷.")]):
    tf = textbox(s, ML + i * (gw + 0.30), 5.24, gw, 0.8)
    para(tf, t, 11.5, BONE, bold=True, spc=0.8, line=1.2, first=True)
    para(tf, b, 9.5, BONE_MID, line=1.5, before=5)

note_card(s, ML, 5.90, SW - ML - MR, "VFX 스코프",
          "동일 세트 · 동일 카메라 위치에서 원혼 덩어리의 볼륨만 단계적으로 확장한다.\n"
          "실물 의상과 배우 연기를 베이스로 덩어리만 CG로 증식시키는 구조 — "
          "촬영 회차를 늘리지 않고 스케일을 확보하는 설계다.",
          size=10, alpha=64, minh=0.86)
foot(s)


# ── 19 · 초기 컨셉 ─────────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "early_white", SW * 0.44, 0, SW * 0.56, SH, focus=0.5)
rect(s, 0, 0, SW * 0.48, SH, INK)
scrim(s, "scrim_left", SW * 0.40, 0, SW * 0.28, SH)
head(s, "17", "ORIGIN OF THE DESIGN", "디자인의 출발점", width=5.3)

tf = textbox(s, ML, 2.20, 5.0, 1.4)
para(tf, "초기 컨셉은 ‘업힌 원혼’ 하나로 시작했다.\n"
         "흰 소복 · 붉은 부적 · 등에 매달린 사람들.",
     14, BONE, bold=True, line=1.45, first=True)
para(tf, "이후 개발에서 세 가지가 확정되었다 — ① 흰 소복을 색동저고리로 바꿔\n"
         "‘지키는 옷’이라는 반전 단서를 심었고, ② 업힌 원혼을 개체에서 유동하는\n"
         "덩어리로 바꿔 하중을 시각화했으며, ③ 절룩이는 다리와 짚신으로\n"
         "그 하중을 걸음걸이에 새겼다.",
     10, BONE_MID, line=1.65, before=14)

note_card(s, ML, 5.20, 5.0, "IP 관점",
          "부적 · 색동 · 짚신 세 기호는 실루엣만으로 식별된다.\n"
          "포스터 · 굿즈 · 게임 · 어트랙션까지 확장 가능한 캐릭터 자산이다.",
          size=9.8, alpha=66)
foot(s)


# ── 20 · 컬러 팔레트 & 키워드 ──────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "hall_sweep", 0, 0, SW, SH, focus=0.62)
scrim(s, "scrim_85")
scrim(s, "vignette")
head(s, "18", "PALETTE & KEYWORDS", "컬러 팔레트 · 컨셉 키워드")

sw_w, gap = 1.24, 0.20
x = ML
for hexv, name in PALETTE:
    rect(s, x, 2.20, sw_w, 1.24, RGBColor.from_string(hexv))
    tf = textbox(s, x, 3.54, sw_w, 0.5)
    para(tf, "#" + hexv, 8, BONE_DIM, font=LAT, spc=0.6, line=1.0, first=True)
    para(tf, name, 9.5, BONE_MID, line=1.2, before=3)
    x += sw_w + gap

hairline(s, ML, 4.52, SW - ML - MR, BONE, 20)

tf = textbox(s, ML, 4.86, 3.2, 0.3)
para(tf, "컨셉 키워드", 9, RED, bold=True, spc=1.8, line=1.0, first=True)
tf = textbox(s, ML, 5.20, 11.4, 0.7)
para(tf, "원한   ·   저주   ·   무거움   ·   억압   ·   끌려다님   ·   구속   ·   공포",
     20, BONE, bold=True, spc=0.6, line=1.25, first=True)

tf = textbox(s, ML, 6.14, 11.4, 0.6)
para(tf, "무채색 위에 단 두 개의 채도 — 색동의 ‘청’과 부적의 ‘홍’. "
         "화면에서 색을 가진 것은 「그것」뿐이며, 그 색은 모두 ‘지키기 위한 색’이다.",
     10.5, BONE_MID, line=1.5, first=True)
foot(s)


# ── 21 · 계보 타임라인 ─────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "corridor_hospital", 0, 0, SW, SH, focus=0.5)
scrim(s, "scrim_85")
scrim(s, "vignette")
head(s, "19", "BLOODLINE", "핏줄의 계보 — 「그것」이 따라온 길")

tl = [("1582", "조선 · 평양", "처형장의 망나니.\n피와 원망이 땅에 쌓인다."),
      ("1587", "굿판", "무당이 노모의 시신을\n부적으로 봉인한다."),
      ("1945", "일제강점기", "순사 백동필.\n제복을 입고 동족을 벤다."),
      ("1984", "일본 오사카", "사형집행관 백종문.\n버튼으로 사람을 죽인다."),
      ("2011", "일본 후쿠오카", "산부인과 의사 슈.\n생명을 지우며 산다."),
      ("현재", "그리고 하루마", "핏줄이 끊긴 자리에\n저주가 되돌아온다.")]
cw2 = (SW - ML - MR) / 6
for i, (yv, place, body) in enumerate(tl):
    x = ML + i * cw2
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(2.42), In(0.10), In(0.10))
    dot.fill.solid(); dot.fill.fore_color.rgb = RED if i in (1, 5) else BONE_MID
    dot.line.fill.background(); dot.shadow.inherit = False
    tf = textbox(s, x, 2.72, cw2 - 0.30, 1.9)
    para(tf, yv, 19, BONE, font=LAT, bold=True, line=1.0, first=True)
    para(tf, place, 9.5, RED_DIM, bold=True, spc=0.8, line=1.2, before=5)
    para(tf, body, 9.3, BONE_MID, line=1.55, before=7)
hairline(s, ML + 0.05, 2.465, SW - ML - MR - cw2 + 0.05, BONE, 22)

note_card(s, ML, 5.28, SW - ML - MR, "왜 이 계보가 중요한가",
          "네 세대 모두 ‘사람을 죽이는 직업’을 가졌다. 칼 → 제복 → 버튼 → 메스. "
          "도구는 문명화되었지만 행위는 같다. 「그것」은 그 사실을 400년 동안 "
          "한 번도 놓치지 않고 지켜본 유일한 목격자다.",
          size=10.5, alpha=60, minh=0.95)
foot(s)


# ── 22 · KEY ART ───────────────────────────────────────────────────────────
s = blank(prs)
nx, ny, nw, nh = pic_fit(s, "keyart", 0, 0, SW * 0.52, SH)
scrim(s, "scrim_right", SW * 0.34, 0, SW * 0.30, SH)
tf = textbox(s, SW * 0.60, 2.05, 4.6, 0.3)
para(tf, "KEY ART", 10, RED, font=LAT, bold=True, spc=3.2, line=1.0, first=True)
tf = textbox(s, SW * 0.60, 2.52, 4.9, 2.6)
para(tf, "罪 · 血 · 記憶", 30, BONE, bold=True, spc=2.0, line=1.25, first=True)
para(tf, "죄 · 피 · 기억", 12, BONE_DIM, spc=2.4, line=1.2, before=6)
para(tf, "티저 키아트는 「그것」을 정면으로 보여주지 않는다.\n"
         "짚신, 색동 자락, 부적, 그리고 벽에 새겨진 세 글자.\n"
         "관객이 조립하게 만드는 방식으로 캐릭터를 노출한다.",
     11, BONE_MID, line=1.65, before=18)
tf = textbox(s, SW * 0.60, 5.60, 4.9, 0.8)
para(tf, "赫血  KAKKETSU", 15, BONE, bold=True, spc=1.6, line=1.2, first=True)
para(tf, "각본 · 감독  한동하   |   제작  WALLAH PROJECT + SUPERNOVA STUDIOS",
     9.5, BONE_DIM, line=1.4, before=6)
foot(s, x=SW * 0.60)


# ── 23 · CLOSING ───────────────────────────────────────────────────────────
s = blank(prs)
pic_cover(s, "kneel_closeup", 0, 0, SW, SH, focus=0.5)
scrim(s, "scrim_85")
scrim(s, "vignette")

tf = textbox(s, ML, 2.30, 9.2, 2.2)
para(tf, "“저주는 사라진 것이 아니다.\n다만, 형태를 바꿨을 뿐이다.”", 28, BONE,
     bold=True, line=1.45, spc=-0.5, first=True)
hairline(s, ML, 4.28, 2.0, BONE, 35)

nxt = [("캐릭터", "룩 C(색동) 기준 확정 · 실물 의상 제작 · 배우 캐스팅"),
       ("VFX", "원혼 덩어리 시뮬레이션 파일럿 (대청 시퀀스 30초)"),
       ("자료", "일본 투자사 제출용 버전 — 문화 주석 확장판 (8월 말)")]
x = ML
for k, v in nxt:
    tf = textbox(s, x, 4.62, 3.6, 0.9)
    para(tf, k, 9, RED, bold=True, spc=1.8, line=1.0, first=True)
    para(tf, v, 10, BONE_MID, line=1.5, before=6)
    x += 3.85

hairline(s, ML, 5.98, SW - ML - MR, BONE, 16)
tf = textbox(s, ML, 6.20, 8.0, 0.6)
para(tf, "赫血  KAKKETSU", 13, BONE, bold=True, spc=1.6, line=1.2, first=True)
para(tf, "WALLAH PROJECT  ×  SUPERNOVA STUDIOS   |   각본 · 감독  한동하",
     9.5, BONE_DIM, line=1.4, before=4)
tf = textbox(s, SW - MR - 4.0, 6.20, 4.0, 0.5)
para(tf, "CONFIDENTIAL", 9, RGBColor(0x6E, 0x66, 0x5D), font=LAT, bold=True,
     spc=2.4, line=1.0, align=PP_ALIGN.RIGHT, first=True)
foot(s)


os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("saved:", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
