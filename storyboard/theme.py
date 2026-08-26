"""
赫血 스토리보드 — 컷 장의 옷을 입힌다.

컷 시트는 「밝은 종이 + 검은 그림 창」이 정석이다.
카드까지 어둡게 하면 어두운 컷 이미지와 뭉개져 지면이 탁해지고,
글은 대비를 아무리 줘도 눅눅하게 읽힌다.
그림은 검은 창 안에 들어가 있어서 종이 위에서 오히려 더 또렷하다.

컷 카드용 표는 슬라이드가 아니라 **레이아웃**에 그려져 있으므로
레이아웃을 한 번만 고치면 모든 컷 장에 적용된다.
"""

from __future__ import annotations

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

PALETTES = {
    "paper": dict(
        bg_top=RGBColor(0xEA, 0xE6, 0xE1), bg_bottom=RGBColor(0xDF, 0xDA, 0xD4),
        plate=RGBColor(0xFF, 0xFF, 0xFF),        # 카드 — 칸을 나누지 않는 흰 판
        label_cell=RGBColor(0xFF, 0xFF, 0xFF),
        head_cell=RGBColor(0xFF, 0xFF, 0xFF),
        well=RGBColor(0x0F, 0x0E, 0x0D),         # 그림이 앉는 창
        empty_well=RGBColor(0xEE, 0xEB, 0xE7),   # 아직 그림이 없는 창
        hairline=RGBColor(0xE2, 0xDD, 0xD7),     # 글 줄 사이의 아주 옅은 선
        body=RGBColor(0x1C, 0x19, 0x17),
        label=RGBColor(0xA5, 0x9D, 0x96),        # 항목 이름 — 배경으로 물러난다
        dim=RGBColor(0xAE, 0xA6, 0x9F),
    ),
    "ink": dict(
        bg_top=RGBColor(0x10, 0x0D, 0x0C), bg_bottom=RGBColor(0x05, 0x05, 0x05),
        plate=RGBColor(0x16, 0x13, 0x11),
        label_cell=RGBColor(0x16, 0x13, 0x11),
        head_cell=RGBColor(0x16, 0x13, 0x11),
        well=RGBColor(0x07, 0x07, 0x06),
        empty_well=RGBColor(0x0B, 0x0A, 0x09),
        hairline=RGBColor(0x3A, 0x32, 0x2D),
        body=RGBColor(0xDC, 0xD7, 0xD2),
        label=RGBColor(0x94, 0x8C, 0x86),
        dim=RGBColor(0x5E, 0x57, 0x52),
    ),
}

_P = PALETTES["paper"]


def use(name: str) -> None:
    """팔레트를 고른다 — paper(기본) 또는 ink."""
    global _P
    _P = PALETTES[name]


def c(key: str) -> RGBColor:
    return _P[key]


# ── 판면 ───────────────────────────────────────────────────────────────
# 2.39:1 컷은 그림 창의 **폭**에서 크기가 정해진다. 창을 세로로 늘려봐야
# 검은 띠만 두꺼워질 뿐이다. 그래서 바깥 여백과 기둥 사이를 걷어내
# 카드 자체를 넓혔다 — 3.788 → 4.158in, 그림이 그만큼 커진다.
SLIDE_W, SLIDE_H = 13.3333, 7.5
MARGIN_X, GUTTER_X = 0.190, 0.240
MARGIN_Y, GUTTER_Y = 0.130, 0.130
CARD_W = (SLIDE_W - MARGIN_X * 2 - GUTTER_X * 2) / 3   # 4.158
CARD_H = (SLIDE_H - MARGIN_Y * 2 - GUTTER_Y) / 2       # 3.555

ROW_H = {          # 인치. 합이 CARD_H 와 같아야 한다
    "head":   0.220,
    "pic":    1.944,   # 4.158 / 1.944 = 2.14:1
    "camera": 0.300,   # 2줄 @8.5pt
    "desc":   0.440,   # 3줄 @8.5pt · 2줄 @9.5pt
    "note":   0.440,
    "trans":  0.211,   # 1줄 @9.5pt
}
ROW_ORDER = ["head", "pic", "camera", "desc", "note", "trans"]
LABEL_COL_W = 0.660   # 항목 이름이 들어가는 왼쪽 칸
LABEL_PT = 7.0        # 항목 이름 — 조용히 물러나 있어야 한다
HEAD_LABEL_PT = 7.0   # Scene / Shot

PICTURE_ROW = 1
FIRST_PH_IDX = 10
PH_PER_PANEL = 7

# 표 선은 글 줄 사이에만 남긴다. 칸마다 테두리를 두르면 엑셀처럼 보인다.
RULED_ROWS = {3, 4, 5}   # 장면묘사 · 디테일 · 전환효과의 윗선


def _set_cell(cell, fill: RGBColor, text_color: RGBColor, size_pt=None) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.text_frame.word_wrap = False   # 항목 이름이 두 줄로 접히지 않게
    for para in cell.text_frame.paragraphs:
        para.font.color.rgb = text_color
        if size_pt:
            para.font.size = Pt(size_pt)
        for run in para.runs:
            run.font.color.rgb = text_color
            if size_pt:
                run.font.size = Pt(size_pt)


def _set_borders(cell, sides: str = "", color: RGBColor | None = None,
                 width_pt: float = 0.5) -> None:
    """
    sides 에 들어간 변만 그린다. 비우면 테두리를 전부 지운다.
    칸마다 사방을 두르면 엑셀처럼 보인다 — 글 줄 사이 윗선 하나면 충분하다.
    """
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    # lnL·lnR·lnT·lnB 는 tcPr 안에서 이 순서로 와야 한다
    for key, tag in (("L", "a:lnL"), ("R", "a:lnR"), ("T", "a:lnT"), ("B", "a:lnB")):
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set("w", str(Pt(width_pt)))
        ln.set("cap", "flat")
        if key in sides and color is not None:
            solid = etree.SubElement(ln, qn("a:solidFill"))
            etree.SubElement(solid, qn("a:srgbClr")).set("val", f"{color}")
        else:
            etree.SubElement(ln, qn("a:noFill"))


def _kill_style_banding(table) -> None:
    """표 스타일이 머리행·줄무늬를 다시 칠하지 못하게 끈다."""
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for attr in ("firstRow", "bandRow", "firstCol", "bandCol", "lastRow", "lastCol"):
            tblPr.attrib.pop(attr, None)


def _rename_cell(cell, text: str) -> None:
    """첫 런만 바꾸면 나머지 런이 남아 「디테일 / 메모」가 된다."""
    paras = cell.text_frame.paragraphs
    first = None
    for para in paras:
        for run in list(para.runs):
            if first is None:
                first, run.text = run, text
            else:
                run._r.getparent().remove(run._r)
    for para in paras[1:]:
        para._p.getparent().remove(para._p)


def _place(ph, left, top, width, height) -> None:
    if ph is None:
        return
    ph.left, ph.top, ph.width, ph.height = int(left), int(top), int(width), int(height)


def restyle_layout(layout) -> dict:
    """
    레이아웃의 컷 카드 6개를 다시 칠하고 칸 비율을 다시 잡는다.
    표만 고치면 글 칸(자리표시자)이 따라오지 않으므로 둘 다 옮긴다.

    돌려주는 값:
      wells       패널마다의 그림 자리 (left, top, width, height)
      boxes       글 칸 하나의 실측 크기 {필드: (너비pt, 높이pt)}
      pic_aspect  그림 창의 가로세로비
    """
    wells, boxes = [], {}
    tables = [s for s in layout.shapes if s.has_table]
    phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}

    for panel, shape in enumerate(tables):
        table = shape.table
        _kill_style_banding(table)

        # 카드를 새 판면 자리로 옮기고 넓힌다 (3열 × 2행, 행 우선)
        col_i, row_i = panel % 3, panel // 3
        shape.left = Inches(MARGIN_X + col_i * (CARD_W + GUTTER_X))
        shape.top = Inches(MARGIN_Y + row_i * (CARD_H + GUTTER_Y))
        shape.width = Inches(CARD_W)
        for row in table.rows:
            if "디테일" in row.cells[0].text:
                _rename_cell(row.cells[0], "디테일")

        # 열 — 왼쪽 항목 칸만 좁히고 나머지를 셋으로 나눈다
        rest = int((shape.width - Inches(LABEL_COL_W)) / 3)
        widths = [Inches(LABEL_COL_W), rest, rest,
                  shape.width - Inches(LABEL_COL_W) - rest * 2]
        for col, w in zip(table.columns, widths):
            col.width = int(w)
        for row, key in zip(table.rows, ROW_ORDER):
            row.height = Inches(ROW_H[key])

        for r, row in enumerate(table.rows):
            for col_i, cell in enumerate(row.cells):
                if r == PICTURE_ROW:
                    fill, text, size = c("well"), c("dim"), None
                elif r == 0:
                    fill, text, size = c("head_cell"), c("label"), HEAD_LABEL_PT
                elif col_i == 0:
                    fill, text, size = c("label_cell"), c("label"), LABEL_PT
                else:
                    fill, text, size = c("plate"), c("body"), None
                _set_cell(cell, fill, text, size)
                _set_borders(cell, "T" if r in RULED_ROWS else "", c("hairline"))
                cell.margin_left = cell.margin_right = Pt(4)
                cell.margin_top = cell.margin_bottom = Pt(1)

        # ── 자리표시자를 새 칸에 맞춰 옮긴다 ──────────────────────────
        base = FIRST_PH_IDX + panel * PH_PER_PANEL
        left, top = shape.left, shape.top
        col0, col1, col2, col3 = widths
        pad = Pt(3)

        head_h, pic_h = Inches(ROW_H["head"]), Inches(ROW_H["pic"])
        y = top + head_h
        _place(phs.get(base), left, y, shape.width, pic_h)
        wells.append((left, y, shape.width, pic_h))

        # Scene 번호 = 두 번째 열, Shot 번호 = 네 번째 열
        _place(phs.get(base + 1), left + col0 + pad, top, col1 - pad * 2, head_h)
        _place(phs.get(base + 2), left + col0 + col1 + col2 + pad, top,
               col3 - pad * 2, head_h)

        y = top + head_h + pic_h
        value_w = col1 + col2 + col3
        for n, key in enumerate(("camera", "desc", "note", "trans"), start=3):
            h = Inches(ROW_H[key])
            _place(phs.get(base + n), left + col0 + pad, y, value_w - pad * 2, h)
            boxes[key] = (Emu(value_w - pad * 2).pt, Emu(h).pt)
            y += h

    for ph in layout.placeholders:
        for para in ph.text_frame.paragraphs:
            para.font.color.rgb = c("body")
            for run in para.runs:
                run.font.color.rgb = c("body")

    return {"wells": wells, "boxes": boxes,
            "pic_aspect": CARD_W / ROW_H["pic"]}


def paint_background(slide) -> None:
    """슬라이드 바탕. 카드보다 한 톤 낮아야 카드가 종이 위에 떠 보인다."""
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    stops = fill.gradient_stops
    stops[0].color.rgb = c("bg_top")
    stops[0].position = 0.0
    stops[1].color.rgb = c("bg_bottom")
    stops[1].position = 1.0


def mark_empty_well(slide, well, label: str = "시안 필요") -> None:
    """
    이미지가 아직 없는 칸. 그림 자리표시자를 그대로 두면
    파워포인트가 「그림을 추가하려면 아이콘을 클릭하세요」를 띄운다.
    대신 옅은 판 하나와 한 줄짜리 안내만 남긴다.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    left, top, width, height = well
    plate = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    plate.fill.solid()
    plate.fill.fore_color.rgb = c("empty_well")
    plate.line.color.rgb = c("hairline")
    plate.line.width = Pt(0.5)
    plate.shadow.inherit = False

    tf = plate.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.color.rgb = c("dim")
    run.font._rPr.set("spc", "200")
