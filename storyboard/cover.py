"""
赫血 스토리보드 — 맨 앞 두 장(표지 · 키 아트)을 만든다.

원본 템플릿의 레이아웃에는 컷 카드용 표 6개가 그려져 있다.
그대로 슬라이드를 만들면 표지 뒤로 그 표가 비쳐 나오므로,
두 장 모두 화면 전체를 덮는 판을 먼저 깔고 그 위에 올린다.
"""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

TITLE_RED = RGBColor(0xA9, 0x00, 0x08)  # 타이틀 로고에서 뽑은 붉은색
BLACK = RGBColor(0x00, 0x00, 0x00)
BONE = RGBColor(0xD8, 0xD4, 0xD0)
ASH = RGBColor(0x8A, 0x85, 0x82)

SLIDE_W = Inches(13.3333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs):
    """레이아웃의 자리표시자를 걷어낸 빈 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    return slide


def _move_to_front(prs, count: int) -> None:
    """마지막에 붙은 슬라이드 count장을 맨 앞으로 옮긴다."""
    id_list = prs.slides._sldIdLst
    ids = list(id_list)
    for i, el in enumerate(ids[-count:]):
        id_list.remove(el)
        id_list.insert(i, el)


def _rect(slide, x, y, w, h, color: RGBColor, alpha: float = 1.0):
    """alpha 1.0 = 불투명. 그보다 작으면 그만큼 비쳐 보인다."""
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if alpha < 1.0:
        clr = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        etree.SubElement(clr, qn("a:alpha")).set("val", str(int(alpha * 100_000)))
    return shape


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(글자, 크기pt, 색, 굵게, 자간pt), ...] — 한 줄에 하나씩."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, spacing) in enumerate(runs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if spacing:
            run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def _fit_into(cell, img_size):
    """셀 안에 비율을 지킨 채 가운데로 넣을 위치·크기."""
    cx, cy, cw, ch = cell
    iw, ih = img_size
    scale = min(cw / iw, ch / ih)
    w, h = int(iw * scale), int(ih * scale)
    return cx + (cw - w) // 2, cy + (ch - h) // 2, w, h


def add_cover(prs, art_dir: Path, deck_cuts: int, total_cuts: int) -> int:
    """표지 · 키 아트 두 장을 맨 앞에 붙이고, 붙인 장수를 돌려준다."""
    from PIL import Image

    made = 0

    # ── 1장: 표지 — 키 아트를 화면 가득 ─────────────────────────────
    cover_img = art_dir / "cover.jpg"
    if cover_img.exists():
        slide = _blank_slide(prs)
        slide.shapes.add_picture(str(cover_img), 0, 0, SLIDE_W, SLIDE_H)

        # 아래쪽을 눌러 글자가 뜨도록 검은 판을 반투명하게 겹친다
        _rect(slide, 0, Inches(5.9), SLIDE_W, Inches(1.6), BLACK, alpha=0.38)

        _rect(slide, Inches(0.75), Inches(6.32), Inches(1.05), Emu(9525), TITLE_RED)

        _text(
            slide, Inches(0.75), Inches(6.52), Inches(7.0), Inches(0.8),
            [
                ("S T O R Y B O A R D", 15, BONE, True, 3.2),
                (f"컷 {deck_cuts}장 수록 · 전체 {total_cuts}컷", 9.5, ASH, False, 0.6),
            ],
        )
        _text(
            slide, Inches(8.0), Inches(6.52), Inches(4.58), Inches(0.8),
            [
                ("각본 · 감독   한 동 하", 11, BONE, False, 1.4),
                ("脚本・監督  ハン・ドンハ    2026.08.22", 8.5, ASH, False, 0.5),
            ],
            align=PP_ALIGN.RIGHT,
        )
        made += 1

    # ── 2장: 키 아트 ────────────────────────────────────────────────
    sheet = [
        ("poster.jpg", (Inches(0.55), Inches(1.15), Inches(4.30), Inches(5.74))),
        ("key_zashiki.jpg", (Inches(5.15), Inches(1.15), Inches(3.70), Inches(2.77))),
        ("key_corridor.jpg", (Inches(9.08), Inches(1.15), Inches(3.70), Inches(2.77))),
        ("key_zashiki2.jpg", (Inches(5.15), Inches(4.12), Inches(3.70), Inches(2.77))),
        ("logo.jpg", (Inches(9.08), Inches(4.12), Inches(3.70), Inches(2.77))),
    ]
    present = [(n, c) for n, c in sheet if (art_dir / n).exists()]
    if present:
        slide = _blank_slide(prs)
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLACK)
        _rect(slide, Inches(0.55), Inches(0.72), Inches(0.9), Emu(9525), TITLE_RED)
        _text(
            slide, Inches(0.55), Inches(0.34), Inches(6.0), Inches(0.4),
            [("K E Y   A R T", 13, BONE, True, 3.0)],
        )
        _text(
            slide, Inches(7.0), Inches(0.34), Inches(5.78), Inches(0.4),
            [("赫血  ·  KAKKETSU", 12, ASH, False, 1.6)],
            align=PP_ALIGN.RIGHT,
        )
        for name, cell in present:
            path = art_dir / name
            with Image.open(path) as im:
                size = im.size
            x, y, w, h = _fit_into(cell, size)
            slide.shapes.add_picture(str(path), x, y, w, h)
        made += 1

    if made:
        _move_to_front(prs, made)
    return made
