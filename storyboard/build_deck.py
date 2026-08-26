#!/usr/bin/env python3
"""
赫血 스토리보드 — 컷 목록(JSON)을 원본 PPT 템플릿에 부어 넣는다.

원본 템플릿은 슬라이드 1장에 패널 6개(3열 × 2행).
패널 하나 = 그림 개체 틀 1 + 텍스트 개체 틀 6
          (Scene / Shot / 카메라 / 장면묘사 / 디테일·메모 / 전환효과)

사용법:
    python build_deck.py                      # cuts.json -> out/storyboard.pptx
    python build_deck.py --cuts cuts.json --template AI_storyboard.pptx \
                         --images images --out out/storyboard.pptx
"""

from __future__ import annotations

import argparse
import copy
import json
import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt

import theme
from cover import add_cover

HERE = Path(__file__).parent

# 만들 때마다 새 이름으로 남긴다. 덮어쓰면 어제 넘긴 판본이 사라진다.
DECK_NAME = "각혈 스토리보드"
SHOOT_NAME = "각혈 촬영컷목록"


def next_version(out_dir: Path, stem: str) -> int:
    """out 폴더에 이미 있는 판본 중 제일 큰 번호 + 1."""
    import re

    found = [
        int(m.group(1))
        for f in out_dir.glob(f"{stem} v*.pptx")
        if (m := re.fullmatch(rf"{re.escape(stem)} v(\d+)", f.stem))
    ]
    return max(found, default=0) + 1

# 패널 1개당 자리표시자 7개가 연속으로 붙어 있다. 패널 6개 × 7 = 42, 시작 idx 10.
PANELS_PER_SLIDE = 6
PH_PER_PANEL = 7
FIRST_PH_IDX = 10

# 텍스트 자리표시자 6개의 순서 (그림 틀 바로 다음부터)
FIELD_ORDER = ["scene", "shot", "camera", "desc", "note", "trans"]

# 칸별 최대 글자 크기— 번호 칸은 크게, 설명 칸은 작게.
# 실제 크기는 칸에 들어가는 만큼 아래에서 자동으로 줄인다.
FIELD_MAX_PT = {
    "scene": 12.0,
    "shot": 12.0,
    "camera": 9.5,
    "desc": 9.5,
    "note": 9.5,
    "trans": 9.0,   # 전환효과는 짧다. 본문보다 앞서 보이면 안 된다
}
# 이보다 작아지면 읽기를 포기하게 된다. 여기 걸리면 글을 줄여야 한다.
FIELD_MIN_PT = 8.0

# 줄간격은 배수가 아니라 pt로 못박는다.
# 배수로 두면 파워포인트가 글꼴 고유 행간(맑은 고딕은 약 1.33em)에 곱해서
# 계산보다 훨씬 두꺼워지고, 아랫칸을 침범한다.
LINE_RATIO = 1.15
# 칸마다 크기가 제각각이면 지면이 지저분해진다. 세 단으로만 떨어뜨린다.
SIZE_STEPS = [9.5, 8.5, 8.0]
NUMBER_STEPS = [12.0, 11.0, 10.0]

# 칸 실측 크기 (너비pt, 높이pt) — theme.restyle_layout 이 채워 넣는다
BOXES: dict[str, tuple[float, float]] = {}


def text_width_units(s: str) -> float:
    """1em을 1.0으로 본 글자 폭의 합. 한글·한자·가나는 1.0, 그 외는 0.5."""
    total = 0.0
    for ch in s:
        total += 1.0 if ord(ch) > 0x2E80 else 0.55
    return total


SAFETY = 0.95   # 글꼴 폭 추정이 빗나가도 아랫칸을 넘지 않도록


def lines_needed(s: str, pt: float, width_pt: float) -> int:
    usable = width_pt * SAFETY
    return max(1, math.ceil(text_width_units(s) * pt / usable))


def fit_font_pt(s: str, field: str) -> tuple[float, bool]:
    """
    칸 안에 다 들어가는 가장 큰 글자 크기.
    두 번째 값은 최소 크기까지 내려도 넘치는지 여부.
    """
    max_pt = FIELD_MAX_PT.get(field, 11.0)
    w, h = BOXES.get(field, (215.0, 24.0))
    usable = h - 2.0
    if not s:
        return max_pt, False
    steps = NUMBER_STEPS if field in ("scene", "shot") else SIZE_STEPS
    for pt in steps:
        if pt > max_pt or pt < FIELD_MIN_PT:
            continue
        if int(usable // (pt * LINE_RATIO)) >= lines_needed(s, pt, w):
            return pt, False
    return FIELD_MIN_PT, True


def fit_to_panel(path: Path, cache: Path, panel_aspect: float) -> Path:
    """
    그림 칸에 넣으면 PowerPoint가 가운데를 기준으로 잘라낸다.
    시네마스코프(2.4:1) 컷은 좌우가 통째로 날아가므로,
    미리 검은 여백을 붙여 칸 비율에 맞춘 사본을 만들어 그걸 넣는다.
    """
    from PIL import Image

    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        aspect = w / h
        if abs(aspect - panel_aspect) < 0.02:
            return path

        if aspect > panel_aspect:      # 더 넓다 -> 위아래에 여백
            nw, nh = w, round(w / panel_aspect)
        else:                          # 더 좁다 -> 좌우에 여백
            nw, nh = round(h * panel_aspect), h

        canvas = Image.new("RGB", (nw, nh), (0, 0, 0))
        canvas.paste(im, ((nw - w) // 2, (nh - h) // 2))

        cache.mkdir(parents=True, exist_ok=True)
        out = cache / f"{path.stem}_{panel_aspect:.3f}.jpg"
        canvas.save(out, quality=92, optimize=True)
        return out


IMAGE_EXTS = [".png", ".jpg", ".jpeg", ".webp"]


def resolve_image(images_dir: Path, name: str) -> Path | None:
    """확장자가 달라도 이름만 맞으면 찾아준다. png로 적어두고 jpg를 넣어도 된다."""
    exact = images_dir / name
    if exact.exists():
        return exact
    stem = Path(name).stem
    for ext in IMAGE_EXTS:
        cand = images_dir / (stem + ext)
        if cand.exists():
            return cand
    return None


def panel_placeholder_indices(panel: int) -> tuple[int, list[int]]:
    """패널 번호(0~5) -> (그림 틀 idx, 텍스트 틀 idx 6개)"""
    base = FIRST_PH_IDX + panel * PH_PER_PANEL
    return base, [base + n for n in range(1, PH_PER_PANEL)]


def clear_slide(slide) -> None:
    """템플릿에 남아 있는 기존 내용(샘플 이미지 포함)을 비운다."""
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def clone_layout_placeholders(slide, layout) -> None:
    """레이아웃의 자리표시자를 슬라이드로 복제한다."""
    spTree = slide.shapes._spTree
    for ph in layout.placeholders:
        spTree.append(copy.deepcopy(ph._element))


def fill_text(ph, value: str, field: str) -> bool:
    value = value or ""
    tf = ph.text_frame
    tf.text = value
    tf.word_wrap = True
    # 칸이 좁으니 안쪽 여백을 최대한 없앤다
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    pt, overflow = fit_font_pt(value, field)
    size = Pt(pt)
    for para in tf.paragraphs:
        para.line_spacing = Pt(pt * LINE_RATIO)
        para.font.size = size
        para.font.color.rgb = theme.c("body")
        for run in para.runs:
            run.font.size = size
            run.font.color.rgb = theme.c("body")
    return overflow


def fill_panel(slide, panel: int, cut: dict, images_dir: Path, cache: Path,
               metrics: dict) -> list[str]:
    pic_idx, text_idxs = panel_placeholder_indices(panel)
    by_idx = {p.placeholder_format.idx: p for p in slide.placeholders}

    tight = []
    for field, idx in zip(FIELD_ORDER, text_idxs):
        ph = by_idx.get(idx)
        if ph is not None and fill_text(ph, str(cut.get(field, "")), field):
            tight.append(field)

    pic_ph = by_idx.get(pic_idx)
    if pic_ph is None:
        return tight

    path = resolve_image(images_dir, cut.get("image") or "")
    if path is not None:
        pic_ph.insert_picture(str(fit_to_panel(path, cache, metrics["pic_aspect"])))
        return tight

    # 아직 그림이 없다 — 자리표시자를 걷어내고 어두운 판만 남긴다
    pic_ph._element.getparent().remove(pic_ph._element)
    wells = metrics["wells"]
    if panel < len(wells):
        theme.mark_empty_well(slide, wells[panel])
    return tight


def drop_unused_placeholders(slide, used_panels: int) -> None:
    """마지막 슬라이드에서 컷이 없는 패널의 자리표시자를 지운다."""
    keep = set()
    for panel in range(used_panels):
        pic_idx, text_idxs = panel_placeholder_indices(panel)
        keep.add(pic_idx)
        keep.update(text_idxs)
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx not in keep:
            ph._element.getparent().remove(ph._element)


def split_cuts(cuts: list[dict], images_dir: Path) -> tuple[list[dict], list[dict]]:
    """
    레퍼런스 덱이므로 그림 칸이 빈 패널을 줄줄이 늘어놓지 않는다.
    이미지가 있는 컷과, 시안이 꼭 있어야 하는 컷(need)만 슬라이드에 싣고
    나머지 — 그대로 촬영하면 되는 컷 — 는 목록으로만 남긴다.
    """
    on_deck, shoot_only = [], []
    for cut in cuts:
        has_image = bool(cut.get("image")) and resolve_image(images_dir, cut["image"])
        (on_deck if has_image or cut.get("need") else shoot_only).append(cut)
    return on_deck, shoot_only


def write_shoot_list(cuts: list[dict], path: Path) -> None:
    """슬라이드에서 뺀 컷을 촬영용 목록으로 남긴다."""
    lines = [
        "# 촬영 컷 목록 — 시안 없이 그대로 찍는 컷",
        "",
        "레퍼런스 덱(`storyboard.pptx`)에서는 뺐지만 컷 번호는 그대로 살아 있습니다.",
        "덱의 씬·컷 번호가 군데군데 건너뛰는 것은 여기 있는 컷들입니다.",
        "",
    ]
    scene = None
    for cut in cuts:
        if cut["scene"] != scene:
            scene = cut["scene"]
            lines += ["", f"## 씬 {scene}", "", "| 컷 | 카메라 | 장면묘사 | 디테일 / 메모 |", "|---|---|---|---|"]
        lines.append(
            f"| {cut['shot']} | {cut['camera']} | {cut['desc']} | {cut.get('note', '')} |"
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build(cuts: list[dict], template: Path, images_dir: Path, out: Path,
          art_dir: Path | None = None, total_cuts: int = 0,
          version: int = 0) -> tuple[int, int]:
    prs = Presentation(str(template))
    layout = prs.slide_layouts[0]

    # 레이아웃을 먼저 고쳐야 한다. 첫 장은 레이아웃을 복제해서 만들기 때문에
    # 순서가 뒤바뀌면 첫 장만 옛 칸 비율을 그대로 물고 온다.
    metrics = theme.restyle_layout(layout)

    # 템플릿에 딸려 온 첫 슬라이드는 비워두고 첫 장으로 재사용한다.
    base_slide = prs.slides[0]
    clear_slide(base_slide)
    clone_layout_placeholders(base_slide, layout)
    BOXES.clear()
    BOXES.update(metrics["boxes"])
    BOXES["scene"] = BOXES["shot"] = (60.0, 20.0)

    slides = [base_slide]
    total_slides = -(-len(cuts) // PANELS_PER_SLIDE)  # ceil
    for _ in range(total_slides - 1):
        s = prs.slides.add_slide(layout)
        slides.append(s)

    tight: list[str] = []
    for i, slide in enumerate(slides):
        chunk = cuts[i * PANELS_PER_SLIDE : (i + 1) * PANELS_PER_SLIDE]
        theme.paint_background(slide)
        for panel, cut in enumerate(chunk):
            for field in fill_panel(slide, panel, cut, images_dir, out.parent / "_fit", metrics):
                tight.append(f"S{cut['scene']}C{cut['shot']} {field}")
        if len(chunk) < PANELS_PER_SLIDE:
            drop_unused_placeholders(slide, len(chunk))

    if tight:
        print(f"  · 8pt로도 칸을 넘치는 글 {len(tight)}곳: {', '.join(tight[:8])}"
              + (" …" if len(tight) > 8 else ""), file=sys.stderr)

    front = add_cover(prs, art_dir, len(cuts), total_cuts, version) if art_dir else 0

    core = prs.core_properties
    core.title = DECK_NAME + (f" v{version:02d}" if version else "")
    core.subject = "赫血 · KAKKETSU"

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return total_slides, front


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--cuts", default=str(HERE / "cuts.json"))
    ap.add_argument("--template", default=str(HERE / "template.pptx"))
    ap.add_argument("--images", default=str(HERE / "images"))
    ap.add_argument("--art", default=str(HERE / "keyart"))
    ap.add_argument("--theme", default="paper", choices=sorted(theme.PALETTES),
                    help="paper = 밝은 종이(기본) · ink = 어두운 판")
    ap.add_argument("--out-dir", default=str(HERE / "out"))
    ap.add_argument("--out", default=None,
                    help="이름을 직접 정하고 싶을 때. 비우면 각혈 스토리보드 v01.pptx 처럼 번호가 붙는다")
    args = ap.parse_args()
    theme.use(args.theme)

    data = json.loads(Path(args.cuts).read_text(encoding="utf-8"))
    cuts = data["cuts"] if isinstance(data, dict) else data

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    if args.out:
        out = Path(args.out)
        version = 0
    else:
        version = next_version(out_dir, DECK_NAME)
        out = out_dir / f"{DECK_NAME} v{version:02d}.pptx"

    on_deck, shoot_only = split_cuts(cuts, Path(args.images))
    shoot_path = out_dir / (f"{SHOOT_NAME} v{version:02d}.md" if version else "촬영컷목록.md")
    write_shoot_list(shoot_only, shoot_path)

    blanks = sum(1 for c in on_deck if not resolve_image(Path(args.images), c.get("image") or ""))
    n, front = build(
        on_deck, Path(args.template), Path(args.images), out,
        art_dir=Path(args.art), total_cuts=len(cuts), version=version,
    )
    print(
        f"전체 {len(cuts)}컷 · 덱 {len(on_deck)}컷(시안 필요 빈칸 {blanks}) · "
        f"촬영 목록 {len(shoot_only)}컷 → 앞장 {front} + 컷 {n} = 슬라이드 {n + front}장"
    )
    print(f"  {out.name}")
    print(f"  {shoot_path.name}")


if __name__ == "__main__":
    main()
